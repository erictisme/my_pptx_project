from dataclasses import dataclass
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json
import logging
import requests
from pathlib import Path

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@dataclass
class SlideContent:
    """Represents the content structure of a slide."""
    title: str
    points: List[str]
    visual_type: str = "bullet_points"
    layout_type: str = "Title and Content"

@dataclass
class StyleGuide:
    """McKinsey style guidelines."""
    fonts: Dict[str, Dict[str, any]] = None
    colors: Dict[str, str] = None
    spacing: Dict[str, float] = None
    
    def __post_init__(self):
        self.fonts = {
            "title": {"name": "Arial", "size": 32},
            "subtitle": {"name": "Arial", "size": 24},
            "body": {"name": "Arial", "size": 18},
            "caption": {"name": "Arial", "size": 14}
        }
        
        self.colors = {
            "primary": "#000000",      # Black
            "secondary": "#0085CA",    # McKinsey Blue
            "background": "#FFFFFF",   # White
            "accent1": "#7B7B7B",     # Gray
            "accent2": "#4A90E2"      # Light Blue
        }
        
        self.spacing = {
            "title_margin_top": 0.5,    # inches
            "content_margin_left": 0.75,
            "paragraph_spacing": 12.0,   # points
            "line_spacing": 1.2         # multiplier
        }

class ContentGenerator:
    """Generates slide content using LLM."""
    
    def __init__(self, model_name: str = "llama2"):
        self.api_base = "http://localhost:11434/api/generate"
        self.model_name = model_name
        
        # Test connection
        try:
            response = requests.get("http://localhost:11434/api/tags")
            response.raise_for_status()
            logger.info("Successfully connected to Ollama")
        except Exception as e:
            logger.error(f"Failed to connect to Ollama: {e}")
            raise
    
    def generate_outline(self, brief: str) -> List[SlideContent]:
        """Generate presentation outline."""
        # Step 1: Generate titles
        titles_prompt = f"""As a McKinsey consultant, create 3 powerful slide titles for this topic:

{brief}

Requirements:
- Action-oriented titles
- Include specific metrics when possible
- Start with verbs
- Focus on business impact

Example:
1. "Drive 45% Cost Reduction Through AI Automation"
2. "Accelerate Growth with Data-Driven Insights"
3. "Capture $2M in Annual Savings"

Your titles:"""

        titles_response = self._generate_text(titles_prompt)
        logger.debug(f"Titles response: {titles_response}")
        
        # Extract titles (assuming numbered list format)
        titles = []
        for line in titles_response.split('\n'):
            line = line.strip()
            if line and (line.startswith('"') or line.startswith('"')):
                title = line.strip('"').strip('"').strip()
                if title:
                    titles.append(title)
        
        if not titles:
            titles = ["AI Business Impact", "Implementation Roadmap", "Expected ROI"]
        
        # Step 2: Generate content for each title
        slides = []
        for title in titles:
            content_prompt = f"""Create 4 powerful bullet points for this slide title:

Title: "{title}"

Requirements:
- Start with action verbs
- Include specific metrics
- Focus on business outcomes
- Keep each point under 8 words

Example:
• Reduce operational costs by 45%
• Increase customer satisfaction to 98%
• Deploy AI solution within 3 months
• Generate $5M additional annual revenue

Your bullet points:"""

            points_response = self._generate_text(content_prompt)
            logger.debug(f"Points response for {title}: {points_response}")
            
            # Extract bullet points
            points = []
            for line in points_response.split('\n'):
                line = line.strip().lstrip('•').strip()
                if line and not line.startswith('Example'):
                    points.append(line)
            
            if not points:
                points = [
                    "Implement solution within 90 days",
                    "Reduce costs by 30%",
                    "Improve efficiency by 40%",
                    "Generate positive ROI in 6 months"
                ]
            
            # Step 3: Determine best visual type
            visual_prompt = f"""What's the best visual type for this slide?

Title: "{title}"
Points:
{chr(10).join('• ' + p for p in points)}

Choose ONE from:
- bar_chart (for metrics over time)
- 2x2_matrix (for comparisons)
- process_flow (for steps)
- timeline (for roadmaps)
- comparison (for before/after)

Answer with just the type:"""

            visual_type = self._generate_text(visual_prompt).strip().lower()
            if visual_type not in ["bar_chart", "2x2_matrix", "process_flow", "timeline", "comparison"]:
                visual_type = "bar_chart"
            
            slides.append(SlideContent(
                title=title,
                points=points[:4],  # Take first 4 points
                visual_type=visual_type
            ))
        
        return slides
    
    def enhance_content(self, slide: SlideContent) -> SlideContent:
        """Enhance slide content with better phrasing and structure."""
        enhance_prompt = f"""Improve these bullet points to be more impactful:

Title: "{slide.title}"
Current points:
{chr(10).join('• ' + p for p in slide.points)}

Requirements:
- Start with strong action verbs
- Include specific metrics
- Focus on business outcomes
- Keep each point under 8 words

Improved points:"""

        response = self._generate_text(enhance_prompt)
        logger.debug(f"Enhancement response: {response}")
        
        # Extract enhanced points
        enhanced_points = []
        for line in response.split('\n'):
            line = line.strip().lstrip('•').strip()
            if line:
                enhanced_points.append(line)
        
        if not enhanced_points:
            enhanced_points = slide.points
        
        return SlideContent(
            title=slide.title,
            points=enhanced_points[:4],
            visual_type=slide.visual_type
        )
    
    def _generate_text(self, prompt: str) -> str:
        """Generate text using Ollama API."""
        data = {
            "model": self.model_name,
            "prompt": f"""You are a McKinsey presentation expert. Be clear and concise.

{prompt}""",
            "stream": False,
            "options": {
                "temperature": 0.7,
                "top_p": 0.9,
                "stop": ["\n\n", "```"]
            }
        }
        
        response = requests.post(self.api_base, json=data)
        response.raise_for_status()
        return response.json()["response"].strip()

class SlideDesigner:
    """Applies McKinsey-style design to slides."""
    
    def __init__(self, style_guide: StyleGuide):
        self.style = style_guide
        self._layout_map = {
            "text_only": "Title and Content",
            "2x2_matrix": "2 columns",
            "process_flow": "Title and Content",
            "bar_chart": "Title and Content",
            "comparison": "2 columns",
            "timeline": "Title and Content"
        }
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    
    def create_slide(self, prs: Presentation, content: SlideContent) -> None:
        """Create a slide with McKinsey styling."""
        # Select layout
        layout_name = self._layout_map.get(content.visual_type, "Title and Content")
        layout = self._get_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(layout)
        
        # Apply title with blue accent bar
        if slide.shapes.title:
            # Add blue accent bar
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(0.1)
            height = Inches(0.5)
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(self.style.colors["secondary"])
            shape.line.fill.background()
            
            # Style title
            title = slide.shapes.title
            title.text = content.title
            title.left = Inches(0.7)  # Adjust for accent bar
            self._apply_text_style(title.text_frame, "title")
        
        # Apply content based on visual type
        if content.visual_type == "2x2_matrix":
            self._create_2x2_matrix(slide, content)
        elif content.visual_type == "bar_chart":
            self._create_chart(slide, content)
        else:
            self._create_bullet_points(slide, content)
        
        # Add footer
        self._add_footer(slide)
    
    def _create_bullet_points(self, slide, content):
        """Create formatted bullet points."""
        body_shape = None
        for shape in slide.placeholders:
            if shape.name in ["Content Placeholder", "Text Placeholder"]:
                body_shape = shape
                break
        
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            
            for point in content.points:
                p = tf.add_paragraph()
                p.text = "•  " + point  # Add bullet with spacing
                self._apply_text_style(p, "body")
                p.space_after = Pt(self.style.spacing["paragraph_spacing"])
                p.space_before = Pt(self.style.spacing["paragraph_spacing"])
    
    def _create_2x2_matrix(self, slide, content):
        """Create a 2x2 matrix layout."""
        # Add matrix box
        left = Inches(2)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        # Draw matrix outline
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
        )
        box.fill.background()
        box.line.color.rgb = self._hex_to_rgb(self.style.colors["accent1"])
        
        # Add dividing lines using thin rectangles
        h_line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left,
            top + height/2 - Inches(0.005),
            width,
            Inches(0.01)
        )
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = self._hex_to_rgb(self.style.colors["accent1"])
        h_line.line.fill.background()
        
        v_line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left + width/2 - Inches(0.005),
            top,
            Inches(0.01),
            height
        )
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = self._hex_to_rgb(self.style.colors["accent1"])
        v_line.line.fill.background()
        
        # Add content to quadrants
        for i, point in enumerate(content.points[:4]):  # Use first 4 points
            quad_left = left + (width/2 * (i % 2))
            quad_top = top + (height/2 * (i // 2))
            
            textbox = slide.shapes.add_textbox(
                quad_left + Inches(0.1),
                quad_top + Inches(0.1),
                width/2 - Inches(0.2),
                height/2 - Inches(0.2)
            )
            
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = point
            self._apply_text_style(p, "body")
            p.alignment = PP_ALIGN.CENTER
    
    def _create_chart(self, slide, content):
        """Create a basic chart."""
        # For now, just create bullet points
        # TODO: Implement actual chart creation
        self._create_bullet_points(slide, content)
    
    def _add_footer(self, slide):
        """Add McKinsey-style footer."""
        footer_shape = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6.5),
            Inches(9),
            Inches(0.5)
        )
        
        tf = footer_shape.text_frame
        p = tf.add_paragraph()
        p.text = "CONFIDENTIAL"
        self._apply_text_style(p, "caption")
        p.alignment = PP_ALIGN.RIGHT
        
        # Add thin line above footer using a thin rectangle instead
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.5),
            Inches(6.45),
            Inches(9),
            Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.style.colors["accent1"])
        line.line.fill.background()
    
    def _get_layout_by_name(self, prs: Presentation, name: str) -> Optional[any]:
        """Get slide layout by name."""
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
        return prs.slide_layouts[1]  # Default to Title and Content
    
    def _apply_text_style(self, text_obj, style_type: str) -> None:
        """Apply text styling to either a text frame or paragraph."""
        font_style = self.style.fonts[style_type]
        
        if hasattr(text_obj, 'paragraphs'):  # It's a text frame
            paragraphs = text_obj.paragraphs
        else:  # It's a paragraph
            paragraphs = [text_obj]
        
        for paragraph in paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = self.style.spacing["line_spacing"]
            
            if not paragraph.runs:
                run = paragraph.add_run()
                run.text = paragraph.text
                paragraph.text = ""
            
            for run in paragraph.runs:
                font = run.font
                font.name = font_style["name"]
                font.size = Pt(font_style["size"])
                font.color.rgb = self._hex_to_rgb(self.style.colors["primary"])

class DeckBuilder:
    """Main class for building McKinsey-style presentations."""
    
    def __init__(self, model_name: str = "llama2", template_path: str = "base_template.pptx"):
        self.style_guide = StyleGuide()
        self.content_generator = ContentGenerator(model_name)
        self.designer = SlideDesigner(self.style_guide)
        self.template_path = template_path
    
    def create_presentation(self, brief: str, output_path: str) -> None:
        """Create a complete presentation."""
        logger.info("Starting presentation creation...")
        
        # Create presentation from template
        prs = Presentation(self.template_path)
        
        try:
            # Generate outline with more specific prompt
            slides = self.content_generator.generate_outline(brief)
            logger.info(f"Generated outline with {len(slides)} slides")
            
            # Create each slide
            for slide in slides:
                # Enhance content with business focus
                enhanced = self.content_generator.enhance_content(slide)
                logger.info(f"Enhanced content for slide: {enhanced.title}")
                
                # Create and design slide
                self.designer.create_slide(prs, enhanced)
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
        except Exception as e:
            logger.error(f"Failed to create presentation: {e}")
            raise

def main():
    """Example usage."""
    brief = """
    Create an executive presentation on AI's business impact:
    - Current value creation
    - Implementation roadmap
    - Expected ROI
    """
    
    try:
        builder = DeckBuilder(template_path="base_template.pptx")
        builder.create_presentation(brief, "mckinsey_ai_impact.pptx")
    except Exception as e:
        logger.error(f"Failed to create presentation: {e}")
        raise

if __name__ == "__main__":
    main() 