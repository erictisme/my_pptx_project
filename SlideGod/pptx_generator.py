import logging
from typing import Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models import StyleGuide

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def create_presentation(style_guide: StyleGuide) -> Presentation:
    """Create a new PowerPoint presentation with the specified style guide."""
    prs = Presentation()
    
    # Set slide size and layout
    prs.slide_width = Inches(style_guide.slide_width)
    prs.slide_height = Inches(style_guide.slide_height)
    
    # Set default font
    for shape in prs.slide_master.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = style_guide.font_family
    
    return prs

def add_slide(prs: Presentation, slide_content: Dict, style_guide: StyleGuide) -> None:
    """Add a new slide to the presentation with the specified content and style."""
    slide_layout = prs.slide_layouts.get_by_name(slide_content["layout"]["name"])
    slide = prs.slides.add_slide(slide_layout)
    
    for shape_data in slide_content["shapes"]:
        shape_type = shape_data["type"]
        location = shape_data["location"]
        text_content = shape_data["textContent"]
        
        if shape_type == "TITLE":
            title_shape = slide.shapes.title
            if not title_shape:
                title_shape = slide.shapes.add_textbox(
                    Inches(location["x"] / 10000000),
                    Inches(location["y"] / 10000000),
                    Inches(location["width"] / 10000000),
                    Inches(location["height"] / 10000000)
                )
            
            title_shape.text = text_content[0]["text"]
            apply_text_style(title_shape.text_frame.paragraphs[0], text_content[0]["style"], style_guide)
        
        elif shape_type == "BODY":
            body_shape = slide.shapes.placeholders[1]
            if not body_shape:
                body_shape = slide.shapes.add_textbox(
                    Inches(location["x"] / 10000000),
                    Inches(location["y"] / 10000000),
                    Inches(location["width"] / 10000000),
                    Inches(location["height"] / 10000000)
                )
            
            tf = body_shape.text_frame
            tf.clear()
            
            for text_data in text_content:
                p = tf.add_paragraph()
                p.text = text_data["text"]
                apply_text_style(p, text_data["style"], style_guide)
        
        else:
            logger.warning(f"Unsupported shape type: {shape_type}")

def apply_text_style(paragraph, style_data: Dict, style_guide: StyleGuide) -> None:
    """Apply text style to a paragraph."""
    font = paragraph.font
    font.name = style_guide.font_family
    font.size = Pt(style_data["fontSize"])
    font.bold = style_data["isBold"]
    font.color.rgb = RGBColor(*style_guide.font_color)
    
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.line_spacing = 1.0
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)

def save_presentation(prs: Presentation, filename: str) -> None:
    """Save the presentation to a file."""
    prs.save(filename)
    logger.info(f"Presentation saved to {filename}")

def generate_presentation(outline: List[Dict], style_guide: StyleGuide, filename: str) -> None:
    """Generate a PowerPoint presentation from the outline."""
    prs = create_presentation(style_guide)
    
    for slide_info in outline:
        add_slide(prs, slide_info, style_guide)
    
    save_presentation(prs, filename) 