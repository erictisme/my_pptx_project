import logging
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class SimplePPTGenerator:
    def __init__(self, model_name="llama2"):
        self.api_base = "http://localhost:11434/api/generate"
        self.model_name = model_name
        
        # Simple styling defaults
        self.style = {
            "colors": {
                "text": RGBColor(0, 0, 0),  # Black
            },
            "fonts": {
                "title": {"name": "Arial", "size": Pt(28)},
                "body": {"name": "Arial", "size": Pt(18)}
            }
        }
    
    def _generate_text(self, prompt: str) -> str:
        """Generate text using Ollama API."""
        data = {
            "model": self.model_name,
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": 0.7,
                "top_p": 0.9,
                "stop": ["\n\n", "```"]
            }
        }
        
        try:
            response = requests.post(self.api_base, json=data)
            response.raise_for_status()
            return response.json()["response"].strip()
        except Exception as e:
            logger.error(f"Error generating text: {e}")
            return ""

    def _create_slide(self, prs: Presentation, title: str, content: str) -> None:
        """Create a slide using Layout 1 (Title and Content)."""
        # Add slide with Title and Content layout
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = self.style["fonts"]["title"]["name"]
                run.font.size = self.style["fonts"]["title"]["size"]
                run.font.color.rgb = self.style["colors"]["text"]
        
        # Add content
        body_shape = slide.placeholders[1]  # Index 1 is the content placeholder
        text_frame = body_shape.text_frame
        text_frame.clear()  # Clear any default text
        
        # Add content as paragraphs
        p = text_frame.add_paragraph()
        p.text = content
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = self.style["fonts"]["body"]["name"]
            run.font.size = self.style["fonts"]["body"]["size"]
            run.font.color.rgb = self.style["colors"]["text"]
    
    def generate(self, brief: str, output_file: str) -> None:
        """Generate a single slide with action-oriented title and content."""
        # Create presentation
        prs = Presentation()
        
        # Generate action-oriented title
        title_prompt = f"""You are a McKinsey consultant creating slide titles.

Topic: {brief}

Create ONE action-oriented slide title following these rules:
1. Start with a STRONG action verb (e.g., Unlock, Capture, Drive, Accelerate)
2. Focus on business outcome, not technology
3. Include numbers when possible (%, $M, 2X)
4. Use connecting words (through, by, via) to link action to method
5. Max 8 words
6. NO QUOTES in response

Strong examples:
- Unlock $100M Value Through Digital Transformation
- Drive 2X Growth by Optimizing Customer Journey
- Accelerate Margin Expansion via Smart Operations

Weak examples (don't use):
- Optimize AI Implementation
- Digital Strategy Overview
- Improve Business with Technology

Your title:"""
        
        title = self._generate_text(title_prompt)
        if not title:
            title = "Unlock 30% Value Through AI-Powered Operations"
        
        # Generate structured content
        content_prompt = f"""Write the main message for this slide:

Topic: {title}

Requirements:
- Start with a clear assertion
- Support with 2-3 specific points
- Include metrics where relevant
- Keep under 4 lines total
- Make it sound like a McKinsey insight
- NO bullet points or special characters

Example:
AI implementation will drive 40% cost reduction across operations, primarily through automated quality control (15% savings) and predictive maintenance (25% savings).

Your content:"""
        
        content = self._generate_text(content_prompt)
        if not content:
            content = "AI implementation can drive 30% efficiency gains across operations, with primary impact in automated processes (20%) and decision support (10%)."
        
        # Create the slide
        self._create_slide(prs, title, content)
        
        # Save presentation
        prs.save(output_file)
        logger.info(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    brief = """
    Create a slide about growth strategy:
    - Market expansion opportunities
    - Customer acquisition
    - New product development
    """
    
    generator = SimplePPTGenerator()
    generator.generate(brief, "simple_mckinsey_slide.pptx") 