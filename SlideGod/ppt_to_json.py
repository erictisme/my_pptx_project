from pptx import Presentation
import json
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

def get_text_style(paragraph):
    """Extract comprehensive text styling information."""
    run = paragraph.runs[0] if paragraph.runs else None
    if not run:
        return {}
    
    style = {
        "fontSize": run.font.size.pt if run.font.size else None,
        "fontName": run.font.name if run.font.name else None,
        "bold": run.font.bold if hasattr(run.font, 'bold') else None,
        "italic": run.font.italic if hasattr(run.font, 'italic') else None,
        "underline": run.font.underline if hasattr(run.font, 'underline') else None,
        "color": f"#{run.font.color.rgb:06x}" if run.font.color and hasattr(run.font.color, 'rgb') else None,
        "alignment": str(paragraph.alignment) if paragraph.alignment else "LEFT",
        "lineSpacing": paragraph.line_spacing if hasattr(paragraph, 'line_spacing') else None,
        "spaceBefore": paragraph.space_before.pt if paragraph.space_before else None,
        "spaceAfter": paragraph.space_after.pt if paragraph.space_after else None
    }
    return {k: v for k, v in style.items() if v is not None}

def get_shape_data(shape):
    """Extract comprehensive shape information."""
    shape_data = {
        "type": str(shape.shape_type),
        "name": shape.name,
        "id": shape.shape_id,
        "location": {
            "x": shape.left,
            "y": shape.top,
            "width": shape.width,
            "height": shape.height,
            "rotation": shape.rotation
        }
    }
    
    # Handle text content
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            para_data = {
                "text": para.text,
                "level": para.level,
                "style": get_text_style(para)
            }
            paragraphs.append(para_data)
        shape_data["textContent"] = paragraphs
        shape_data["verticalAnchor"] = str(shape.text_frame.vertical_anchor)
        shape_data["wordWrap"] = shape.text_frame.word_wrap
        
    # Handle tables
    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                cell_data = {
                    "text": cell.text,
                    "style": get_text_style(cell.text_frame.paragraphs[0])
                }
                row_data.append(cell_data)
            table_data.append(row_data)
        shape_data["tableContent"] = table_data
    
    return shape_data

def extract_slide_data(pptx_file):
    """Extract comprehensive presentation data."""
    prs = Presentation(pptx_file)
    presentation_data = {
        "slides": [],
        "metadata": {
            "slideWidth": prs.slide_width,
            "slideHeight": prs.slide_height,
            "slideMasterStyles": []  # TODO: Extract master styles
        }
    }
    
    for slide in prs.slides:
        slide_data = {
            "layout": {
                "name": slide.slide_layout.name,
                "type": slide.slide_layout.slide_master.name
            },
            "shapes": [],
            "background": {
                "fill": "TODO: Extract background fill"  # TODO: Implement background extraction
            }
        }
        
        for shape in slide.shapes:
            shape_data = get_shape_data(shape)
            slide_data["shapes"].append(shape_data)
        
        presentation_data["slides"].append(slide_data)
    
    return presentation_data

if __name__ == "__main__":
    pptx_file = "Leading in Permacrisis Workshop Singapore Dec 6.pptx"
    data = extract_slide_data(pptx_file)
    
    # Save with pretty printing for readability
    with open("enhanced_slides.json", "w", encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    
    print("Enhanced conversion complete! Check enhanced_slides.json")
