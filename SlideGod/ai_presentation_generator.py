import json
import os
import logging
from typing import List, Dict, Any, Set
import requests
from abc import ABC, abstractmethod
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from dataclasses import dataclass
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

from models import (
    SlideContent,
    StyleGuide,
    ResearchInput,
    ProjectObjective,
    ProjectContext
)
from api_calls import (
    generate_outline,
    generate_slide_content,
    enhance_content
)
from pptx_utils import (
    hex_to_rgb,
    apply_text_style,
    create_shape_from_data,
    create_slide_from_data
)
from config import (
    TEMPLATE_FILE,
    OUTPUT_DIR,
    LOGGING_LEVEL,
    LOGGING_FORMAT
)
from pptx_generator import generate_presentation

# Set up logging
logging.basicConfig(
    level=LOGGING_LEVEL,
    format=LOGGING_FORMAT
)
logger = logging.getLogger(__name__)

class PresentationError(Exception):
    """Custom exception for presentation-related errors."""
    pass

class StyleExtractionError(PresentationError):
    """Error during style pattern extraction."""
    pass

class ContentGenerationError(PresentationError):
    """Error during content generation."""
    pass

class PresentationStyleExtractor:
    """Extracts and learns styling patterns from example presentations."""
    
    def __init__(self):
        self.style_patterns = {
            "colors": {},
            "fonts": {},
            "spacing": {},
            "layouts": {}
        }
    
    def analyze_presentation(self, pptx_file: str):
        """Analyze a presentation to extract styling patterns."""
        prs = Presentation(pptx_file)
        
        # Analyze color schemes
        self._extract_color_scheme(prs)
        
        # Analyze font usage
        self._extract_font_patterns(prs)
        
        # Analyze spacing and alignment
        self._extract_spacing_patterns(prs)
        
        # Analyze layout usage
        self._extract_layout_patterns(prs)
    
    def _extract_color_scheme(self, prs: Presentation):
        """Extract color scheme from the presentation."""
        color_usage = {
            "background": {},
            "text": {},
            "accent": {}
        }
        
        for slide in prs.slides:
            # Extract background colors safely
            try:
                if hasattr(slide, 'background') and slide.background.fill.type != 0:  # 0 is no fill
                    if hasattr(slide.background.fill, 'fore_color') and hasattr(slide.background.fill.fore_color, 'rgb'):
                        color = f"#{slide.background.fill.fore_color.rgb:06x}"
                        color_usage["background"][color] = color_usage["background"].get(color, 0) + 1
            except (AttributeError, TypeError):
                pass
            
            # Extract text and shape colors
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run.font, 'color') and run.font.color and hasattr(run.font.color, 'rgb'):
                                    color = f"#{run.font.color.rgb:06x}"
                                    color_usage["text"][color] = color_usage["text"].get(color, 0) + 1
                
                    # Extract shape fill colors
                    if hasattr(shape, 'fill') and shape.fill.type != 0:  # Not a no-fill
                        if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                            color = f"#{shape.fill.fore_color.rgb:06x}"
                            color_usage["accent"][color] = color_usage["accent"].get(color, 0) + 1
                except (AttributeError, TypeError):
                    continue
        
        # Find most common colors
        self.style_patterns["colors"] = {
            "primary": max(color_usage["text"].items(), key=lambda x: x[1])[0] if color_usage["text"] else "#000000",
            "background": max(color_usage["background"].items(), key=lambda x: x[1])[0] if color_usage["background"] else "#FFFFFF",
            "accent": [k for k, v in sorted(color_usage["accent"].items(), key=lambda x: x[1], reverse=True)[:3]] if color_usage["accent"] else []
        }
    
    def _extract_font_patterns(self, prs: Presentation):
        """Extract font usage patterns."""
        font_usage = {
            "names": {},
            "sizes": {
                "title": [],
                "body": [],
                "header": []
            },
            "styles": {
                "bold": 0,
                "italic": 0,
                "underline": 0
            }
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Extract font names
                        if run.font.name:
                            font_usage["names"][run.font.name] = font_usage["names"].get(run.font.name, 0) + 1
                        
                        # Extract font sizes
                        if run.font.size:
                            size = run.font.size.pt
                            if shape.is_placeholder:
                                if shape.placeholder_format.type == 1:  # Title
                                    font_usage["sizes"]["title"].append(size)
                                elif shape.placeholder_format.type == 2:  # Body
                                    font_usage["sizes"]["body"].append(size)
                                elif shape.placeholder_format.type == 3:  # Header
                                    font_usage["sizes"]["header"].append(size)
                        
                        # Extract font styles
                        if run.font.bold:
                            font_usage["styles"]["bold"] += 1
                        if run.font.italic:
                            font_usage["styles"]["italic"] += 1
                        if run.font.underline:
                            font_usage["styles"]["underline"] += 1
        
        # Process collected data
        self.style_patterns["fonts"] = {
            "primary": max(font_usage["names"].items(), key=lambda x: x[1])[0] if font_usage["names"] else "Arial",
            "sizes": {
                "title": round(sum(font_usage["sizes"]["title"]) / len(font_usage["sizes"]["title"])) if font_usage["sizes"]["title"] else 32,
                "body": round(sum(font_usage["sizes"]["body"]) / len(font_usage["sizes"]["body"])) if font_usage["sizes"]["body"] else 18,
                "header": round(sum(font_usage["sizes"]["header"]) / len(font_usage["sizes"]["header"])) if font_usage["sizes"]["header"] else 24
            },
            "styles": {
                k: v > len(prs.slides) * 0.5 for k, v in font_usage["styles"].items()
            }
        }
    
    def _extract_spacing_patterns(self, prs: Presentation):
        """Extract spacing and alignment patterns."""
        spacing_data = {
            "paragraph": {
                "before": [],
                "after": [],
                "line": []
            },
            "margins": {
                "left": [],
                "right": [],
                "top": [],
                "bottom": []
            },
            "alignment": {
                "left": 0,
                "center": 0,
                "right": 0,
                "justify": 0
            }
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Collect shape margins
                spacing_data["margins"]["left"].append(shape.left)
                spacing_data["margins"]["top"].append(shape.top)
                spacing_data["margins"]["right"].append(prs.slide_width - (shape.left + shape.width))
                spacing_data["margins"]["bottom"].append(prs.slide_height - (shape.top + shape.height))
                
                for paragraph in shape.text_frame.paragraphs:
                    # Collect paragraph spacing
                    if paragraph.space_before:
                        spacing_data["paragraph"]["before"].append(paragraph.space_before.pt)
                    if paragraph.space_after:
                        spacing_data["paragraph"]["after"].append(paragraph.space_after.pt)
                    if paragraph.line_spacing:
                        spacing_data["paragraph"]["line"].append(paragraph.line_spacing)
                    
                    # Collect alignment
                    if paragraph.alignment:
                        if paragraph.alignment == PP_ALIGN.LEFT:
                            spacing_data["alignment"]["left"] += 1
                        elif paragraph.alignment == PP_ALIGN.CENTER:
                            spacing_data["alignment"]["center"] += 1
                        elif paragraph.alignment == PP_ALIGN.RIGHT:
                            spacing_data["alignment"]["right"] += 1
                        elif paragraph.alignment == PP_ALIGN.JUSTIFY:
                            spacing_data["alignment"]["justify"] += 1
        
        # Process collected data
        self.style_patterns["spacing"] = {
            "paragraph": {
                "before": round(sum(spacing_data["paragraph"]["before"]) / len(spacing_data["paragraph"]["before"])) if spacing_data["paragraph"]["before"] else 6,
                "after": round(sum(spacing_data["paragraph"]["after"]) / len(spacing_data["paragraph"]["after"])) if spacing_data["paragraph"]["after"] else 6,
                "line": round(sum(spacing_data["paragraph"]["line"]) / len(spacing_data["paragraph"]["line"]), 1) if spacing_data["paragraph"]["line"] else 1.15
            },
            "margins": {
                k: round(sum(v) / len(v)) if v else 36 for k, v in spacing_data["margins"].items()
            },
            "preferred_alignment": max(spacing_data["alignment"].items(), key=lambda x: x[1])[0]
        }
    
    def _extract_layout_patterns(self, prs: Presentation):
        """Extract layout usage patterns."""
        layout_usage = {}
        layout_elements = {}
        
        for slide in prs.slides:
            layout_name = slide.slide_layout.name
            layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1
            
            if layout_name not in layout_elements:
                layout_elements[layout_name] = {
                    "placeholders": {},
                    "shapes": set(),
                    "typical_elements": set(),
                    "content_structure": {
                        "title_location": None,
                        "body_structure": None,
                        "has_footer": False,
                        "has_header": False,
                        "grid_layout": None
                    }
                }
            
            # Analyze placeholders
            title_shape = None
            body_shapes = []
            
            for shape in slide.placeholders:
                ph_type = str(shape.placeholder_format.type)
                layout_elements[layout_name]["placeholders"][ph_type] = layout_elements[layout_name]["placeholders"].get(ph_type, 0) + 1
                
                # Track title and body locations
                if shape.placeholder_format.type == 1:  # Title
                    title_shape = shape
                elif shape.placeholder_format.type == 2:  # Body
                    body_shapes.append(shape)
            
            # Analyze content structure
            if title_shape:
                layout_elements[layout_name]["content_structure"]["title_location"] = {
                    "top": title_shape.top,
                    "left": title_shape.left,
                    "width": title_shape.width,
                    "height": title_shape.height
                }
            
            # Determine body structure
            if body_shapes:
                # Check if body shapes form a grid
                lefts = sorted(set(shape.left for shape in body_shapes))
                tops = sorted(set(shape.top for shape in body_shapes))
                
                if len(lefts) > 1 and len(tops) > 1:
                    layout_elements[layout_name]["content_structure"]["grid_layout"] = {
                        "columns": len(lefts),
                        "rows": len(tops)
                    }
            
            # Check for header/footer
            for shape in slide.shapes:
                if shape.top < prs.slide_height * 0.1:  # Top 10%
                    layout_elements[layout_name]["content_structure"]["has_header"] = True
                if shape.top > prs.slide_height * 0.9:  # Bottom 10%
                    layout_elements[layout_name]["content_structure"]["has_footer"] = True
                
                # Analyze shapes and typical elements
                if hasattr(shape, 'shape_type'):
                    layout_elements[layout_name]["shapes"].add(str(shape.shape_type))
                
                if shape.has_chart:
                    layout_elements[layout_name]["typical_elements"].add("chart")
                if shape.has_table:
                    layout_elements[layout_name]["typical_elements"].add("table")
                if shape.has_text_frame:
                    layout_elements[layout_name]["typical_elements"].add("text")
        
        # Process collected data
        self.style_patterns["layouts"] = {
            "common_layouts": [k for k, v in sorted(layout_usage.items(), key=lambda x: x[1], reverse=True)],
            "layout_patterns": {
                name: {
                    "frequency": count,
                    "elements": {
                        "placeholders": elements["placeholders"],
                        "shapes": list(elements["shapes"]),
                        "typical_elements": list(elements["typical_elements"])
                    },
                    "content_structure": elements["content_structure"]
                }
                for name, count in layout_usage.items()
                for elements in [layout_elements[name]]
            }
        }
    
    def save_patterns(self, output_file: str):
        """Save extracted patterns to a JSON file."""
        with open(output_file, 'w') as f:
            json.dump(self.style_patterns, f, indent=2)

class BaseContentGenerator(ABC):
    """Abstract base class for content generation."""
    
    @abstractmethod
    def generate_outline(self, brief: str) -> List[Dict[str, Any]]:
        pass
    
    @abstractmethod
    def generate_slide_content(self, slide_info: Dict[str, Any]) -> Dict[str, Any]:
        pass

class OllamaContentGenerator(BaseContentGenerator):
    """Generates presentation content using Ollama models."""
    
    def __init__(self, model_name: str = "llama2"):
        """Initialize with model name (default: llama2)."""
        self.model_name = model_name
        self.api_base = "http://localhost:11434/api/generate"
        
        # Verify Ollama service is available
        try:
            response = requests.get("http://localhost:11434/api/tags")
            response.raise_for_status()
            models = response.json().get("models", [])
            available_models = [m["name"].split(":")[0] for m in models]
            
            if model_name not in available_models:
                logger.warning(f"Model {model_name} not found. Available models: {', '.join(available_models)}")
                if available_models:
                    self.model_name = available_models[0]
                    logger.info(f"Using {self.model_name} as fallback model")
            
            # Configure model parameters
            self.generation_config = {
                "context_length": 4096,  # Match training context length
                "num_ctx": 4096,
                "temperature": 0.7,
                "top_p": 0.9,
                "top_k": 40,
                "num_predict": 1024,  # Reasonable limit for slide content
                "stop": ["\n\n", "```"],  # Stop at clear boundaries
                "repeat_penalty": 1.1  # Slight penalty for repetition
            }
            
            logger.info(f"Successfully connected to Ollama service with model: {self.model_name}")
            logger.info(f"Using context length: {self.generation_config['context_length']}")
        except requests.exceptions.ConnectionError:
            logger.error("Could not connect to Ollama service. Please ensure Ollama is running.")
            raise ContentGenerationError("Ollama service not available")
        except Exception as e:
            logger.error(f"Error initializing Ollama service: {str(e)}")
            raise ContentGenerationError(f"Failed to initialize Ollama service: {str(e)}")
    
    def _extract_json_from_response(self, response: str) -> str:
        """Extract JSON from LLM response, handling various formats."""
        logger.debug("Original response:")
        logger.debug(response)
        
        # Remove any markdown code block indicators and surrounding whitespace
        clean_response = response.replace('```json', '').replace('```', '').strip()
        
        # If we don't have a JSON array yet, try to generate one from the key points
        if not (clean_response.startswith('[') or clean_response.startswith('{')):
            # Extract lines that look like bullet points
            lines = clean_response.split('\n')
            points = []
            for line in lines:
                if line.strip().startswith(('-', '*', '•')):
                    points.append(line.strip().lstrip('-*• ').strip())
            
            if points:
                # Create a basic JSON structure from the points
                json_structure = [
                    {
                        "title": "Key Points",
                        "type": "executive_summary",
                        "key_points": points,
                        "visuals": "Simple bullet point layout"
                    }
                ]
                return json.dumps(json_structure, indent=2)
        
        # If the response doesn't start with [ or {, try to find them
        if not clean_response.startswith('[') and not clean_response.startswith('{'):
            json_start = clean_response.find('[')
            if json_start == -1:
                json_start = clean_response.find('{')
            if json_start == -1:
                raise ContentGenerationError("No JSON structure found in response")
            clean_response = clean_response[json_start:]
        
        # If the response doesn't end with ] or }, try to find them
        if not clean_response.endswith(']') and not clean_response.endswith('}'):
            json_end = clean_response.rfind(']')
            if json_end == -1:
                json_end = clean_response.rfind('}')
            if json_end == -1:
                raise ContentGenerationError("No JSON structure found in response")
            clean_response = clean_response[:json_end + 1]
        
        logger.debug("Extracted JSON structure:")
        logger.debug(clean_response)
        
        # Replace Python-style values
        replacements = [
            ("'", '"'),  # Replace single quotes with double quotes
            ('True', 'true'),
            ('False', 'false'),
            ('None', 'null')
        ]
        
        for old, new in replacements:
            clean_response = clean_response.replace(old, new)
        
        # Remove any trailing commas in arrays and objects
        lines = clean_response.split('\n')
        clean_lines = []
        for i, line in enumerate(lines):
            line = line.rstrip()
            if not line:
                continue
            
            # Remove trailing commas before closing brackets
            if line.rstrip().endswith(','):
                next_line = next((l.strip() for l in lines[i + 1:] if l.strip()), '')
                if next_line.startswith('}') or next_line.startswith(']'):
                    line = line.rstrip(',')
            
            clean_lines.append(line)
        
        json_str = '\n'.join(clean_lines)
        logger.debug("Cleaned JSON string:")
        logger.debug(json_str)
        
        # Validate JSON structure
        try:
            # Try to parse it to catch any remaining issues
            parsed = json.loads(json_str)
            # If successful, return the re-serialized (properly formatted) JSON
            return json.dumps(parsed, indent=2)
        except json.JSONDecodeError as e:
            logger.error(f"JSON validation failed: {str(e)}")
            logger.error("Problematic JSON:")
            logger.error(json_str)
            raise ContentGenerationError(f"Failed to validate JSON structure: {str(e)}")

    def _generate_text(self, prompt: str) -> str:
        """Generate text using Ollama API with optimized parameters."""
        data = {
            "model": self.model_name,
            "prompt": f"""You are a professional presentation generator. Your task is to generate content in JSON format.

Rules:
1. Output ONLY valid JSON
2. No markdown, no explanations
3. Keep responses concise and focused

{prompt}""",
            "stream": False,
            "options": self.generation_config
        }
        
        logger.info(f"Sending request to Ollama API with model: {self.model_name}")
        try:
            response = requests.post(self.api_base, json=data, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            if "error" in result:
                raise ContentGenerationError(f"Ollama API error: {result['error']}")
            
            logger.info("Successfully received response from Ollama API")
            return result["response"].strip()
        except requests.exceptions.Timeout:
            error_msg = "Request to Ollama API timed out after 30 seconds"
            logger.error(error_msg)
            raise ContentGenerationError(error_msg)
        except requests.exceptions.ConnectionError:
            error_msg = "Lost connection to Ollama service"
            logger.error(error_msg)
            raise ContentGenerationError(error_msg)
        except Exception as e:
            error_msg = f"Unexpected error while generating content: {str(e)}"
            logger.error(error_msg)
            raise ContentGenerationError(error_msg)
    
    def generate_outline(self, brief: str) -> List[Dict[str, Any]]:
        """Generate a presentation outline from a brief."""
        prompt = f"""Create a McKinsey-style presentation outline.

Brief: {brief}

Example format:
[{{
  "t": "Executive Summary",  // title
  "k": ["Key point 1", "Key point 2"],  // key points
  "v": "2x2 matrix"  // visual type
}}]

Requirements:
1. Each slide needs clear message
2. Use action-oriented titles
3. 3-4 key points per slide
4. Specific visual type

Return array of slides following this exact format."""

        try:
            response = self._generate_text(prompt)
            logger.debug(f"Raw response from model: {response}")
            
            # Extract and clean JSON from response
            json_str = self._extract_json_from_response(response)
            logger.debug(f"Cleaned JSON string: {json_str}")
            
            # Parse and validate JSON structure
            content = json.loads(json_str)
            
            # Convert compact format to full format
            if isinstance(content, list):
                outline = []
                for slide in content:
                    full_slide = {
                        "title": slide.get("t", "Untitled Slide"),
                        "type": self._infer_slide_type(slide.get("t", "")),
                        "key_points": slide.get("k", []),
                        "visuals": slide.get("v", "Basic layout")
                    }
                    outline.append(full_slide)
                return outline
            
            raise ContentGenerationError("Generated content is not a list")
            
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON from model response: {str(e)}")
            return self._generate_fallback_outline(brief)
        except Exception as e:
            logger.error(f"Unexpected error in generate_outline: {str(e)}")
            raise ContentGenerationError(f"Failed to generate outline: {str(e)}")

    def _infer_slide_type(self, title: str) -> str:
        """Infer slide type from title."""
        title_lower = title.lower()
        if "summary" in title_lower or "overview" in title_lower:
            return "executive_summary"
        if "problem" in title_lower or "challenge" in title_lower:
            return "problem_statement"
        if "solution" in title_lower or "approach" in title_lower:
            return "solution"
        if "next" in title_lower or "step" in title_lower:
            return "roadmap"
        if any(word in title_lower for word in ["data", "metric", "number", "stat"]):
            return "data"
        if "conclusion" in title_lower or "recommendation" in title_lower:
            return "conclusion"
        return "content"

    def _generate_fallback_outline(self, brief: str) -> List[Dict[str, Any]]:
        """Generate a basic outline as fallback."""
        return [
            {
                "title": "Executive Summary",
                "type": "executive_summary",
                "key_points": [brief.split('\n')[0]],
                "visuals": "Title slide with key message"
            }
        ]
    
    def generate_slide_content(self, slide_info: Dict[str, Any]) -> Dict[str, Any]:
        """Generate detailed content for a single slide."""
        # Use compact format for prompt
        prompt = f"""Create slide content.

Input: {{
  "t": "{slide_info['title']}",
  "type": "{slide_info['type']}",
  "k": {json.dumps(slide_info['key_points'])},
  "v": "{slide_info['visuals']}"
}}

Return JSON with this structure:
{{
  "l": {{  // layout
    "n": "type_name",
    "t": "content_type"
  }},
  "s": [  // shapes
    {{
      "t": "TITLE",  // type
      "p": {{  // position
        "x": 1000000,
        "y": 1000000,
        "w": 8000000,
        "h": 1000000
      }},
      "c": [  // content
        {{
          "t": "text",
          "s": {{  // style
            "f": 32,  // font size
            "b": true  // bold
          }}
        }}
      ]
    }}
  ]
}}

Rules:
1. Keep JSON compact
2. Use short keys
3. Include all points
4. Match visual type"""

        try:
            response = self._generate_text(prompt)
            logger.debug(f"Raw slide content response: {response}")
            
            # Extract and clean JSON from response
            json_str = self._extract_json_from_response(response)
            logger.debug(f"Cleaned JSON string: {json_str}")
            
            # Parse compact format
            content = json.loads(json_str)
            
            # Convert compact format to full format
            full_content = {
                "layout": {
                    "name": content.get("l", {}).get("n", slide_info["type"].title()),
                    "type": content.get("l", {}).get("t", "content")
                },
                "shapes": []
            }
            
            # Convert shapes to full format
            for shape in content.get("s", []):
                full_shape = {
                    "type": shape.get("t", "BODY"),
                    "location": {
                        "x": shape.get("p", {}).get("x", 1000000),
                        "y": shape.get("p", {}).get("y", 2500000),
                        "width": shape.get("p", {}).get("w", 8000000),
                        "height": shape.get("p", {}).get("h", 4000000)
                    },
                    "textContent": []
                }
                
                # Convert text content
                for text in shape.get("c", []):
                    full_text = {
                        "text": text.get("t", ""),
                        "style": {
                            "fontSize": text.get("s", {}).get("f", 24),
                            "isBold": text.get("s", {}).get("b", False)
                        }
                    }
                    full_shape["textContent"].append(full_text)
                
                full_content["shapes"].append(full_shape)
            
            # Ensure required shapes exist
            if not any(s["type"] == "TITLE" for s in full_content["shapes"]):
                full_content["shapes"].insert(0, {
                    "type": "TITLE",
                    "location": {
                        "x": 1000000,
                        "y": 1000000,
                        "width": 8000000,
                        "height": 1000000
                    },
                    "textContent": [{
                        "text": slide_info["title"],
                        "style": {
                            "fontSize": 32,
                            "isBold": True
                        }
                    }]
                })
            
            if not any(s["type"] == "BODY" for s in full_content["shapes"]):
                full_content["shapes"].append({
                    "type": "BODY",
                    "location": {
                        "x": 1000000,
                        "y": 2500000,
                        "width": 8000000,
                        "height": 4000000
                    },
                    "textContent": [{
                        "text": point,
                        "style": {
                            "fontSize": 24,
                            "isBold": False
                        }
                    } for point in slide_info["key_points"]]
                })
            
            return full_content
            
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON from model response: {str(e)}")
            return self._generate_fallback_slide_content(slide_info)
        except Exception as e:
            logger.error(f"Unexpected error in generate_slide_content: {str(e)}")
            raise ContentGenerationError(f"Failed to generate slide content: {str(e)}")

    def _generate_fallback_slide_content(self, slide_info: Dict[str, Any]) -> Dict[str, Any]:
        """Generate basic slide content as fallback."""
        return {
            "layout": {
                "name": slide_info["type"].title(),
                "type": "basic"
            },
            "shapes": [
                {
                    "type": "TITLE",
                    "location": {
                        "x": 1000000,
                        "y": 1000000,
                        "width": 8000000,
                        "height": 1000000
                    },
                    "textContent": [{
                        "text": slide_info["title"],
                        "style": {
                            "fontSize": 32,
                            "isBold": True
                        }
                    }]
                },
                {
                    "type": "BODY",
                    "location": {
                        "x": 1000000,
                        "y": 2500000,
                        "width": 8000000,
                        "height": 4000000
                    },
                    "textContent": [{
                        "text": point,
                        "style": {
                            "fontSize": 24,
                            "isBold": False
                        }
                    } for point in slide_info["key_points"]]
                }
            ]
        }

class OpenAIContentGenerator(BaseContentGenerator):
    """Original OpenAI-based generator (kept for reference)."""
    
    def __init__(self, api_key: str):
        import openai
        self.api_key = api_key
        openai.api_key = api_key
    
    def generate_outline(self, brief: str) -> List[Dict[str, Any]]:
        """Generate a presentation outline from a brief."""
        # Original OpenAI implementation
        pass
    
    def generate_slide_content(self, slide_info: Dict[str, Any]) -> Dict[str, Any]:
        """Generate detailed content for a single slide."""
        # Original OpenAI implementation
        pass

class PresentationGenerator:
    """Main class for generating McKinsey-style presentations."""
    
    def __init__(self, style_patterns: Dict, model_name: str = "llama2"):
        self.style_patterns = style_patterns
        self.model_name = model_name
    
    def generate_presentation(self, context: ProjectContext, output_file: str):
        """Generate a complete presentation from a project context."""
        logger.info("Starting presentation generation...")
        
        # Generate outline
        outline = generate_outline(context, self.model_name)
        
        # Generate detailed content for each slide
        slides = []
        for slide_info in outline:
            detailed_content = generate_slide_content(slide_info, context, self.model_name)
            enhanced_content = enhance_content(detailed_content, context, self.model_name)
            styled_content = self._apply_styling(enhanced_content)
            slides.append(styled_content)
        
        # Convert to PowerPoint
        presentation_data = {
            "metadata": self.style_patterns,
            "slides": slides
        }
        
        self._create_pptx(presentation_data, TEMPLATE_FILE, output_file)
        logger.info(f"Presentation generated successfully: {output_file}")
    
    def _apply_styling(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Apply McKinsey-style patterns to the content."""
        # Get slide type and find matching layout
        slide_type = content["layout"]["type"]
        layout_patterns = self.style_patterns["layouts"]["layout_patterns"]
        
        # Find best matching layout based on content type and elements
        best_layout = self._find_best_layout(content, layout_patterns)
        
        # Apply the chosen layout
        if best_layout:
            content["layout"]["name"] = best_layout
        
        # Apply styling to shapes
        for shape in content["shapes"]:
            if "textContent" in shape:
                for text_item in shape["textContent"]:
                    style = text_item.get("style", {})
                    
                    # Apply font styling
                    self._apply_font_style(shape["type"], style)
                    
                    # Apply colors
                    style["color"] = self.style_patterns["colors"]["primary"]
                    
                    # Apply spacing
                    style["paragraphSpacing"] = {
                        "before": self.style_patterns["spacing"]["paragraph"]["before"],
                        "after": self.style_patterns["spacing"]["paragraph"]["after"],
                        "line": self.style_patterns["spacing"]["paragraph"]["line"]
                    }
                    
                    # Apply alignment
                    style["alignment"] = self.style_patterns["spacing"]["preferred_alignment"]
                    
                    text_item["style"] = style
        
        return content
    
    def _find_best_layout(self, content: Dict, layout_patterns: Dict) -> str:
        """Find the best matching layout for the content."""
        best_layout = None
        best_score = -1
        
        for layout_name, pattern in layout_patterns.items():
            score = self._score_layout(content, pattern)
            if score > best_score:
                best_score = score
                best_layout = layout_name
        
        return best_layout
    
    def _score_layout(self, content: Dict, pattern: Dict) -> int:
        """Score a layout based on how well it matches the content."""
        score = 0
        
        # Score based on typical elements
        required_elements = self._get_required_elements(content)
        matching_elements = set(pattern["elements"]["typical_elements"]) & required_elements
        score += len(matching_elements) * 2
        
        # Bonus for frequently used layouts
        score += pattern["frequency"]
        
        return score
    
    def _get_required_elements(self, content: Dict) -> Set[str]:
        """Get the required elements based on the content."""
        required_elements = set()
        if "textContent" in str(content):
            required_elements.add("text")
        if "tableContent" in str(content):
            required_elements.add("table")
        if "chartContent" in str(content):
            required_elements.add("chart")
        return required_elements
    
    def _apply_font_style(self, shape_type: str, style: Dict):
        """Apply font styling based on the shape type."""
        if shape_type == "TITLE":
            style["fontSize"] = self.style_patterns["fonts"]["sizes"]["title"]
        elif shape_type == "BODY":
            style["fontSize"] = self.style_patterns["fonts"]["sizes"]["body"]
        else:
            style["fontSize"] = self.style_patterns["fonts"]["sizes"]["body"]
        
        style["fontName"] = self.style_patterns["fonts"]["primary"]
        style["bold"] = self.style_patterns["fonts"]["styles"]["bold"]
        style["italic"] = self.style_patterns["fonts"]["styles"]["italic"]
        style["underline"] = self.style_patterns["fonts"]["styles"]["underline"]
    
    def _create_pptx(self, presentation_data: Dict, template_file: str, output_file: str):
        """Create a PowerPoint presentation from the presentation data."""
        prs = Presentation(template_file)
        
        # Set presentation-level metadata if available
        if "metadata" in presentation_data:
            # TODO: Apply presentation-level styling
            pass
        
        # Create slides
        for slide_data in presentation_data["slides"]:
            create_slide_from_data(prs, slide_data)
        
        # Save the presentation
        output_path = os.path.join(OUTPUT_DIR, output_file)
        prs.save(output_path)
        logger.info(f"Presentation saved to {output_path}")

def load_project_context(filename: str) -> ProjectContext:
    """Load project context from a JSON file."""
    with open(filename, 'r') as f:
        context_data = json.load(f)
    return ProjectContext.from_dict(context_data)

def load_style_guide(filename: str) -> StyleGuide:
    """Load style guide from a JSON file."""
    with open(filename, 'r') as f:
        style_data = json.load(f)
    return StyleGuide.from_dict(style_data)

def main():
    # Load project context and style guide
    context = load_project_context('project_context.json')
    style_guide = load_style_guide('style_guide.json')
    
    # Generate presentation outline
    outline = generate_outline(context, 'llama2')
    logger.info(f"Generated {len(outline)} slides in outline")
    
    # Generate and enhance slide content
    enhanced_outline = []
    for slide_info in outline:
        slide_content = generate_slide_content(slide_info, context, 'llama2')
        enhanced_content = enhance_content(slide_content, context, 'llama2')
        enhanced_outline.append(enhanced_content)
    
    # Generate PowerPoint presentation
    generate_presentation(enhanced_outline, style_guide, 'presentation.pptx')

if __name__ == '__main__':
    main() 