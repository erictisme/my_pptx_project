from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

def hex_to_rgb(hex_color):
    """Convert hex color to RGB."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def apply_text_style(text_frame, style_data):
    """Apply comprehensive text styling to a shape."""
    if not style_data:
        return
    
    paragraph = text_frame.paragraphs[0]
    
    # Apply paragraph-level styling
    if "alignment" in style_data:
        alignment_map = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFY": PP_ALIGN.JUSTIFY
        }
        paragraph.alignment = alignment_map.get(style_data["alignment"], PP_ALIGN.LEFT)
    
    if "lineSpacing" in style_data:
        paragraph.line_spacing = style_data["lineSpacing"]
    
    if "spaceBefore" in style_data:
        paragraph.space_before = Pt(style_data["spaceBefore"])
    
    if "spaceAfter" in style_data:
        paragraph.space_after = Pt(style_data["spaceAfter"])
    
    # Apply run-level styling
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    font = run.font
    
    if "fontSize" in style_data:
        font.size = Pt(style_data["fontSize"])
    
    if "fontName" in style_data:
        font.name = style_data["fontName"]
    
    if "bold" in style_data:
        font.bold = style_data["bold"]
    
    if "italic" in style_data:
        font.italic = style_data["italic"]
    
    if "underline" in style_data:
        font.underline = style_data["underline"]
    
    if "color" in style_data:
        font.color.rgb = hex_to_rgb(style_data["color"])

def create_shape_from_data(slide, shape_data):
    """Create and style a shape based on the provided data."""
    # Handle text shapes
    if "textContent" in shape_data:
        left = Inches(shape_data["location"]["x"] / 914400)  # Convert to inches
        top = Inches(shape_data["location"]["y"] / 914400)
        width = Inches(shape_data["location"]["width"] / 914400)
        height = Inches(shape_data["location"]["height"] / 914400)
        
        shape = slide.shapes.add_textbox(left, top, width, height)
        
        # Clear existing paragraphs
        text_frame = shape.text_frame
        text_frame.clear()
        
        # Add paragraphs from data
        for para_data in shape_data["textContent"]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = para_data["text"]
            apply_text_style(text_frame, para_data.get("style", {}))
    
    # Handle tables
    elif "tableContent" in shape_data:
        rows = len(shape_data["tableContent"])
        cols = len(shape_data["tableContent"][0]) if rows > 0 else 0
        
        left = Inches(shape_data["location"]["x"] / 914400)
        top = Inches(shape_data["location"]["y"] / 914400)
        width = Inches(shape_data["location"]["width"] / 914400)
        height = Inches(shape_data["location"]["height"] / 914400)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        for i, row in enumerate(shape_data["tableContent"]):
            for j, cell_data in enumerate(row):
                cell = table.cell(i, j)
                cell.text = cell_data["text"]
                apply_text_style(cell.text_frame, cell_data.get("style", {}))

def create_slide_from_data(prs, slide_data):
    """Create a slide from the comprehensive slide data."""
    # Find matching layout
    layout_name = slide_data["layout"]["name"]
    matching_layout = None
    
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            matching_layout = layout
            break
    
    if not matching_layout:
        matching_layout = prs.slide_layouts[0]  # Default to first layout
    
    # Create slide
    slide = prs.slides.add_slide(matching_layout)
    
    # Create shapes
    for shape_data in slide_data["shapes"]:
        create_shape_from_data(slide, shape_data)
    
    return slide

def json_to_pptx(json_data, template_file, output_file):
    """Create a PowerPoint presentation from the comprehensive JSON data."""
    prs = Presentation(template_file)
    
    # Set presentation-level metadata if available
    if "metadata" in json_data:
        # TODO: Apply presentation-level styling
        pass
    
    # Create slides
    for slide_data in json_data["slides"]:
        create_slide_from_data(prs, slide_data)
    
    # Save the presentation
    prs.save(output_file)
    print(f"Enhanced presentation saved as {output_file}")

if __name__ == "__main__":
    # Load the enhanced JSON data
    with open("enhanced_slides.json", "r", encoding='utf-8') as f:
        presentation_data = json.load(f)
    
    template_file = "base_template.pptx"
    output_file = "enhanced_output.pptx"
    
    json_to_pptx(presentation_data, template_file, output_file)
